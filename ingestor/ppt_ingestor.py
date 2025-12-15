# ppt_ingestor.py
from pptx import Presentation
from pathlib import Path

class PPTIngestor:
    """
    Extracts text slide-by-slide from PPTX files.
    Returns:
        {
            "slides": [{slide_number, text}],
            "text": full_text,
            "source_type": "ppt"
        }
    """

    @staticmethod
    def parse(file_path: str):
        file_path = Path(file_path)
        if not file_path.exists():
            return {"error": "PPTX file not found"}

        try:
            prs = Presentation(file_path)

            slides_text = []
            full_text = []

            for i, slide in enumerate(prs.slides, start=1):
                slide_text_parts = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        slide_text_parts.append(shape.text)

                slide_text = "\n".join(slide_text_parts).strip()
                full_text.append(slide_text)

                slides_text.append({
                    "slide_number": i,
                    "text": slide_text
                })

            return {
                "slides": slides_text,
                "text": "\n".join(full_text),
                "source_type": "ppt"
            }

        except Exception as e:
            return {"error": f"PPT parse error: {str(e)}"}
